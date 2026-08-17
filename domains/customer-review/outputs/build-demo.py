# Build the 12-page Fuling demo deck for the customer-review domain.
# Pages reuse the same framework/component functions as the component source deck
# (identical structure, themeColor semantics); numbers are the acceptance reference
# values from 客户分析技能包/03_测试数据/参考值_涪陵2026上半年.md.
# NOTE: demo pages are structurally identical to the component source pages but are
# built directly for the demo; formal usage duplicates the source slides via the
# component flow. Data: 涪陵丽桥 2026 H1 (reference values, 2026-08-14).
import os
import sys

sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.abspath(__file__))
MASTER = os.path.join(BASE, '..', 'assets', 'reference-decks', 'customer-review-template.pptx')
OUT = os.path.join(BASE, '涪陵丽桥_客户经营分析汇报_2026上半年_demo.pptx')

import json
CONTENT = json.load(open(os.path.join(BASE, 'demo-content.json'), encoding='utf-8'))


ACC1 = MSO_THEME_COLOR.ACCENT_1
ACC2 = MSO_THEME_COLOR.ACCENT_2
ACC3 = MSO_THEME_COLOR.ACCENT_3
ACC5 = MSO_THEME_COLOR.ACCENT_5
DK1 = MSO_THEME_COLOR.DARK_1
DK2 = MSO_THEME_COLOR.DARK_2
LT2 = MSO_THEME_COLOR.LIGHT_2
FONT = 'Microsoft YaHei'


def add_textbox(slide, name, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, 0)
    return tf


def set_text(tf, text, size, bold, theme_color, align=PP_ALIGN.LEFT, line_spacing=None):
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = FONT
    r.font.color.theme_color = theme_color


def add_rect(slide, name, x, y, w, h, fill_color=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.name = name
    if fill_color is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.theme_color = fill_color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def furniture(slide):
    add_rect(slide, 'TOP_BAR', 0, 0, 13.33, 0.06, fill_color=ACC1)
    tf = add_textbox(slide, 'FOOTER', 0.6, 7.06, 12.1, 0.35)
    set_text(tf, '客户经营分析 · 汇报 ｜ 涪陵丽桥 · 2026 上半年', 10, False, ACC5)


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    furniture(s)
    return s


def header(s, title, subtitle):
    tf = add_textbox(s, 'TITLE', 0.6, 0.42, 11.8, 0.75)
    set_text(tf, title, 26, True, DK1)
    tf = add_textbox(s, 'SUBTITLE', 0.6, 1.08, 11.8, 0.5)
    set_text(tf, subtitle, 14, False, ACC5)


prs = Presentation(MASTER)

# P1 cover
s = new_slide(prs)
tf = add_textbox(s, 'KICKER', 1.0, 2.0, 11.3, 0.5)
set_text(tf, CONTENT['cover']['kicker'], 14, True, ACC1, PP_ALIGN.CENTER)
tf = add_textbox(s, 'COVER_TITLE', 0.8, 2.6, 11.7, 1.2)
set_text(tf, CONTENT['cover']['title'], 38, True, DK1, PP_ALIGN.CENTER)
tf = add_textbox(s, 'COVER_SUB', 0.8, 3.9, 11.7, 0.6)
set_text(tf, CONTENT['cover']['sub'], 18, False, ACC5, PP_ALIGN.CENTER)
tf = add_textbox(s, 'COVER_META', 0.8, 6.0, 11.7, 0.6)
set_text(tf, CONTENT['cover']['meta'], 12, False, ACC5, PP_ALIGN.CENTER)

# P2 one-sentence conclusion
s = new_slide(prs)
header(s, '一句话经营结论', '全盘先看三件事')
concl = CONTENT["conclusion"]
for i, (label, value, note) in enumerate(concl):
    y = 1.9 + i * 1.55
    add_rect(s, f'C{i+1}_BAND', 0.9, y, 11.5, 1.25, fill_color=LT2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = add_textbox(s, f'C{i+1}_LABEL', 1.2, y + 0.2, 3.0, 0.4)
    set_text(tf, label, 13, False, ACC5)
    tf = add_textbox(s, f'C{i+1}_VALUE', 1.2, y + 0.6, 5.0, 0.55)
    set_text(tf, value, 22, True, DK1)
    tf = add_textbox(s, f'C{i+1}_NOTE', 6.5, y + 0.55, 5.6, 0.6)
    set_text(tf, note, 12, False, ACC5)

# P3 KPI overview (component KPI_OVERVIEW)
s = new_slide(prs)
header(s, 'KPI 总览', '核心经营指标（2026 上半年）')
kpis = [('KPI_' + str(i + 1), *row) for i, row in enumerate(CONTENT['kpis'])]
for i, (prefix, label, value, note) in enumerate(kpis):
    x = 0.6 if i % 2 == 0 else 6.9
    y = 1.7 if i < 2 else 4.2
    add_rect(s, prefix + '_BAND', x, y, 5.9, 2.3, fill_color=LT2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = add_textbox(s, prefix + '_LABEL', x + 0.3, y + 0.25, 5.3, 0.4)
    set_text(tf, label, 12, False, ACC5)
    tf = add_textbox(s, prefix + '_VALUE', x + 0.3, y + 0.75, 5.3, 0.85)
    set_text(tf, value, 30, True, DK1)
    tf = add_textbox(s, prefix + '_NOTE', x + 0.3, y + 1.65, 5.3, 0.5)
    set_text(tf, note, 11, False, ACC5)

# P4 tier table (component TIER_TABLE)
s = new_slide(prs)
header(s, '客户分层结果', 'A/B/C/D 四级：客户数与销售额占比')
tiers = CONTENT["tiers"]
cols = [('TIER', 0.6, 1.4), ('COUNT', 2.2, 1.8), ('SALES', 4.2, 2.4), ('SHARE', 6.8, 1.4), ('NOTE', 8.4, 4.0)]
for name, x, w in cols:
    tf = add_textbox(s, 'COL_' + name, x, 1.7, w, 0.45)
    set_text(tf, {'TIER': '层级', 'COUNT': '客户数', 'SALES': '销售额', 'SHARE': '占比', 'NOTE': '经营含义'}[name], 13, True, DK2)
for i, (tier, count, sales, share, note) in enumerate(tiers):
    y = 2.25 + i * 0.85
    if i % 2 == 0:
        add_rect(s, f'TIER_{i+1}_BAND', 0.6, y, 12.1, 0.8, fill_color=LT2)
    for name, x, w in cols:
        val = {'TIER': tier, 'COUNT': count, 'SALES': sales, 'SHARE': share, 'NOTE': note}[name]
        tf = add_textbox(s, f'TIER_{i+1}_{name}', x, y + 0.18, w, 0.5)
        set_text(tf, val, 13, name == 'TIER', ACC1 if name == 'TIER' and i == 0 else DK1)
tf = add_textbox(s, 'TOTAL_COUNT', 2.2, 5.75, 1.8, 0.45)
set_text(tf, '1,329', 13, True, DK1)
tf = add_textbox(s, 'TOTAL_SALES', 4.2, 5.75, 2.4, 0.45)
set_text(tf, '39,698,859', 13, True, DK1)
tf = add_textbox(s, 'CONCLUSION', 0.6, 6.25, 12.1, 0.5)
set_text(tf, CONTENT['tier_conclusion'], 15, True, DK1)

# P5 customer card (component CUSTOMER_CARD, anonymous example)
s = new_slide(prs)
header(s, CONTENT['card']['title'], CONTENT['card']['subtitle'])
fields = [('FIELD_' + str(i + 1), label, value) for i, (label, value) in enumerate(CONTENT['card']['fields'])]
for i, (prefix, label, value) in enumerate(fields):
    x = 0.6 if i < 3 else 6.9
    y = 1.75 + (i % 3) * 0.8
    tf = add_textbox(s, prefix + '_LABEL', x, y, 1.9, 0.5)
    set_text(tf, label, 13, False, ACC5)
    tf = add_textbox(s, prefix + '_VALUE', x + 2.0, y, 4.0, 0.5)
    set_text(tf, value, 13, True, DK1)
tf = add_textbox(s, 'ADVICE_LABEL', 0.6, 4.5, 2.2, 0.5)
set_text(tf, '下一步建议', 14, True, ACC1)
tf = add_textbox(s, 'ADVICE_TEXT', 2.9, 4.5, 9.8, 1.6)
set_text(tf, CONTENT['card']['advice'], 13, False, DK1, line_spacing=1.3)

# P6 tag distribution (KPI grid reuse)
s = new_slide(prs)
header(s, '客户标签分布', '标签计数（2026-06-30 基准）')
tags = [('TAG_' + str(i + 1), *row) for i, row in enumerate(CONTENT['tags'])]
for i, (prefix, label, value, note) in enumerate(tags):
    x = 0.6 if i % 2 == 0 else 6.9
    y = 1.7 if i < 2 else 4.2
    add_rect(s, prefix + '_BAND', x, y, 5.9, 2.3, fill_color=LT2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = add_textbox(s, prefix + '_LABEL', x + 0.3, y + 0.25, 5.3, 0.4)
    set_text(tf, label, 12, False, ACC5)
    tf = add_textbox(s, prefix + '_VALUE', x + 0.3, y + 0.75, 5.3, 0.85)
    set_text(tf, value, 30, True, DK1)
    tf = add_textbox(s, prefix + '_NOTE', x + 0.3, y + 1.65, 5.3, 0.5)
    set_text(tf, note, 11, False, ACC5)

# P7 sleeping activation (component LIST_TABLE)
s = new_slide(prs)
header(s, CONTENT['sleep']['title'], CONTENT['sleep']['subtitle'])
sleep_rows = CONTENT["sleep"]["rows"]
cols = [('NAME', 0.6, 3.4), ('METRIC', 4.2, 2.8), ('STATUS', 7.2, 1.8), ('ACTION', 9.2, 3.5)]
for name, x, w in cols:
    tf = add_textbox(s, 'COL_' + name, x, 1.7, w, 0.45)
    set_text(tf, {'NAME': '沉睡档位', 'METRIC': '客户数', 'STATUS': '状态', 'ACTION': '激活建议'}[name], 13, True, DK2)
for i, (name, metric, status, action) in enumerate(sleep_rows):
    y = 2.25 + i * 0.85
    if i % 2 == 0:
        add_rect(s, f'ROW_{i+1}_BAND', 0.6, y, 12.1, 0.8, fill_color=LT2)
    for name2, x, w in cols:
        val = {'NAME': name, 'METRIC': metric, 'STATUS': status, 'ACTION': action}[name2]
        tf = add_textbox(s, f'ROW_{i+1}_{name2}', x, y + 0.18, w, 0.5)
        color = ACC2 if '深度' in status else DK1
        set_text(tf, val, 13, False, color)
tf = add_textbox(s, 'FOOT_NOTE', 0.6, 5.5, 12.1, 0.5)
set_text(tf, CONTENT['sleep']['note'], 11, False, ACC5)

# P8 new-customer conversion (KPI grid)
s = new_slide(prs)
header(s, CONTENT['conversion']['title'], CONTENT['conversion']['subtitle'])
conv = [('CONV_' + str(i + 1), *row) for i, row in enumerate(CONTENT['conversion']['kpis'])]
for i, (prefix, label, value, note) in enumerate(conv):
    x = 0.6 if i % 2 == 0 else 6.9
    y = 1.7 if i < 2 else 4.2
    add_rect(s, prefix + '_BAND', x, y, 5.9, 2.3, fill_color=LT2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = add_textbox(s, prefix + '_LABEL', x + 0.3, y + 0.25, 5.3, 0.4)
    set_text(tf, label, 12, False, ACC5)
    tf = add_textbox(s, prefix + '_VALUE', x + 0.3, y + 0.75, 5.3, 0.85)
    set_text(tf, value, 26, True, DK1)
    tf = add_textbox(s, prefix + '_NOTE', x + 0.3, y + 1.65, 5.3, 0.5)
    set_text(tf, note, 11, False, ACC5)

# P9 churn warning (component LIST_TABLE)
s = new_slide(prs)
header(s, CONTENT['churn']['title'], CONTENT['churn']['subtitle'])
churn = CONTENT["churn"]["rows"]
for name, x, w in cols:
    tf = add_textbox(s, 'COL_' + name, x, 1.7, w, 0.45)
    set_text(tf, {'NAME': '预警类型', 'METRIC': '客户数', 'STATUS': '判定特征', 'ACTION': '挽回建议'}[name], 13, True, DK2)
for i, (name, metric, status, action) in enumerate(churn):
    y = 2.25 + i * 0.85
    if i % 2 == 0:
        add_rect(s, f'ROW_{i+1}_BAND', 0.6, y, 12.1, 0.8, fill_color=LT2)
    for name2, x, w in cols:
        val = {'NAME': name, 'METRIC': metric, 'STATUS': status, 'ACTION': action}[name2]
        tf = add_textbox(s, f'ROW_{i+1}_{name2}', x, y + 0.18, w, 0.5)
        set_text(tf, val, 13, False, ACC2 if name2 == 'STATUS' else DK1)
tf = add_textbox(s, 'FOOT_NOTE', 0.6, 5.5, 12.1, 0.5)
set_text(tf, CONTENT['churn']['note'], 11, False, ACC5)

# P10 maintenance plan (component LIST_TABLE)
s = new_slide(prs)
header(s, CONTENT['maintenance']['title'], CONTENT['maintenance']['subtitle'])
maint = CONTENT["maintenance"]["rows"]
cols2 = [('NAME', 0.6, 2.8), ('METRIC', 3.6, 3.6), ('STATUS', 7.4, 3.2), ('ACTION', 10.6, 2.2)]
for name, x, w in cols2:
    tf = add_textbox(s, 'COL_' + name, x, 1.7, w, 0.45)
    set_text(tf, {'NAME': '维护项', 'METRIC': '规则', 'STATUS': '说明', 'ACTION': '标注'}[name], 13, True, DK2)
for i, (name, metric, status, action) in enumerate(maint):
    y = 2.25 + i * 0.85
    if i % 2 == 0:
        add_rect(s, f'ROW_{i+1}_BAND', 0.6, y, 12.1, 0.8, fill_color=LT2)
    for name2, x, w in cols2:
        val = {'NAME': name, 'METRIC': metric, 'STATUS': status, 'ACTION': action}[name2]
        tf = add_textbox(s, f'ROW_{i+1}_{name2}', x, y + 0.18, w, 0.5)
        set_text(tf, val, 13, False, DK1)
tf = add_textbox(s, 'FOOT_NOTE', 0.6, 5.5, 12.1, 0.5)
set_text(tf, CONTENT['maintenance']['note'], 11, False, ACC5)

# P11 thirty-day actions
s = new_slide(prs)
header(s, '未来 30 天行动', '对象 · 动作 · 责任人 · 节点 · 验收')
actions = CONTENT["actions"]
for i, (act, obj, action, owner, date, metric) in enumerate(actions):
    y = 1.7 + i * 1.15
    add_rect(s, f'A{i+1}_BAND', 0.6, y, 12.1, 1.0, fill_color=LT2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = add_textbox(s, f'A{i+1}_ACT', 0.9, y + 0.12, 1.6, 0.4)
    set_text(tf, act, 13, True, ACC1)
    tf = add_textbox(s, f'A{i+1}_OBJ', 2.6, y + 0.12, 2.2, 0.4)
    set_text(tf, obj, 13, True, DK1)
    tf = add_textbox(s, f'A{i+1}_TEXT', 0.9, y + 0.55, 8.5, 0.4)
    set_text(tf, action, 12, False, DK1)
    tf = add_textbox(s, f'A{i+1}_OWNER', 9.5, y + 0.15, 1.6, 0.4)
    set_text(tf, owner, 12, False, ACC5)
    tf = add_textbox(s, f'A{i+1}_DATE', 11.2, y + 0.15, 1.4, 0.4)
    set_text(tf, date, 12, False, ACC5)
    tf = add_textbox(s, f'A{i+1}_METRIC', 9.5, y + 0.55, 3.1, 0.4)
    set_text(tf, metric, 12, False, ACC3)

# P12 closing
s = new_slide(prs)
tf = add_textbox(s, 'CLOSE_TITLE', 0.8, 2.4, 11.7, 0.9)
set_text(tf, '下一步', 30, True, DK1, PP_ALIGN.CENTER)
tf = add_textbox(s, 'CLOSE_TEXT', 1.2, 3.6, 10.9, 2.0)
set_text(tf, CONTENT['closing'], 16, False, DK1, PP_ALIGN.LEFT, line_spacing=1.5)

prs.save(OUT)
print('demo written:', OUT, 'slides:', len(prs.slides))
