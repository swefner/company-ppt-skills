# Build the four generic customer-review report components on the domain master.
# All semantic colors use themeColor references (accent slots), so the whole set
# recolors when the theme changes. Page furniture (top bar + footer) is added per slide.
import os
import sys

sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.abspath(__file__))
MASTER = os.path.join(BASE, '..', 'reference-decks', 'customer-review-template.pptx')
OUT = os.path.join(BASE, 'customer-review-components.pptx')

ACC1 = MSO_THEME_COLOR.ACCENT_1
ACC2 = MSO_THEME_COLOR.ACCENT_2
ACC3 = MSO_THEME_COLOR.ACCENT_3
ACC4 = MSO_THEME_COLOR.ACCENT_4
ACC5 = MSO_THEME_COLOR.ACCENT_5
ACC6 = MSO_THEME_COLOR.ACCENT_6
DK1 = MSO_THEME_COLOR.DARK_1
DK2 = MSO_THEME_COLOR.DARK_2
LT2 = MSO_THEME_COLOR.LIGHT_2

FONT = 'Microsoft YaHei'


def add_textbox(slide, name, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def set_text(tf, text, size, bold, theme_color, align=PP_ALIGN.LEFT, line_spacing=None):
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = FONT
    r.font.color.theme_color = theme_color


def add_rect(slide, name, x, y, w, h, fill_color=None, line_color=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.name = name
    if fill_color is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.theme_color = fill_color
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.theme_color = line_color
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def add_furniture(slide):
    add_rect(slide, 'TOP_BAR', 0, 0, 13.33, 0.06, fill_color=ACC1)
    tf = add_textbox(slide, 'FOOTER', 0.6, 7.06, 12.1, 0.35)
    set_text(tf, '客户经营分析 · 汇报', 10, False, ACC5)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_furniture(slide)
    return slide


def header(slide, title, subtitle):
    tf = add_textbox(slide, 'TITLE', 0.6, 0.42, 11.8, 0.75)
    set_text(tf, title, 26, True, DK1)
    tf = add_textbox(slide, 'SUBTITLE', 0.6, 1.08, 11.8, 0.5)
    set_text(tf, subtitle, 14, False, ACC5)


def cue(slide, text):
    tf = add_textbox(slide, 'TEACHING_CUE', 0.6, 6.6, 12.1, 0.5)
    set_text(tf, text, 11, False, ACC5)


prs = Presentation(MASTER)

# ---- Component 1: KPI Overview (KPI_OVERVIEW) ----
s = new_slide(prs)
header(s, 'KPI 总览', '本期间核心经营指标')
kpis = [
    ('KPI_1', 0.6, 1.7), ('KPI_2', 6.9, 1.7), ('KPI_3', 0.6, 4.2), ('KPI_4', 6.9, 4.2),
]
for prefix, x, y in kpis:
    add_rect(s, prefix + '_BAND', x, y, 5.9, 2.3, fill_color=LT2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = add_textbox(s, prefix + '_LABEL', x + 0.3, y + 0.25, 5.3, 0.4)
    set_text(tf, '指标名称', 12, False, ACC5)
    tf = add_textbox(s, prefix + '_VALUE', x + 0.3, y + 0.75, 5.3, 0.85)
    set_text(tf, '0.00', 30, True, DK1)
    tf = add_textbox(s, prefix + '_NOTE', x + 0.3, y + 1.65, 5.3, 0.5)
    set_text(tf, '口径说明', 11, False, ACC5)
cue(s, '讲解提示：逐块说明指标口径与变化含义；重点指标可用主题强调色标注。')

# ---- Component 2: Tier Table (TIER_TABLE) ----
s = new_slide(prs)
header(s, '客户分层结果', 'A/B/C/D 四级：客户数、销售额与占比')
cols = [('COL_TIER', 0.6, 1.4), ('COL_COUNT', 2.2, 1.8), ('COL_SALES', 4.2, 2.4),
        ('COL_SHARE', 6.8, 1.4), ('COL_NOTE', 8.4, 4.0)]
for name, x, w in cols:
    tf = add_textbox(s, name, x, 1.7, w, 0.45)
    set_text(tf, {'COL_TIER': '层级', 'COL_COUNT': '客户数', 'COL_SALES': '销售额',
                  'COL_SHARE': '占比', 'COL_NOTE': '经营含义'}[name], 13, True, DK2)
for i in range(1, 5):
    y = 2.25 + (i - 1) * 0.85
    fill = ACC1 if i == 1 else None
    if fill is not None:
        # A-tier row uses a light accent1 tint via XML not supported by python-pptx;
        # use LT2 band for all rows, emphasize tier-1 name with accent1 text.
        fill = None
    if fill:
        add_rect(s, f'TIER_{i}_BAND', 0.6, y, 12.1, 0.8, fill_color=fill)
    for name, x, w in cols:
        key = name.replace('COL_', '')
        tf = add_textbox(s, f'TIER_{i}_{key}', x, y + 0.15, w, 0.5)
        color = ACC1 if (key == 'TIER' and i == 1) else DK1
        set_text(tf, ['A', 'B', 'C', 'D'][i - 1] if key == 'TIER' else '—', 13,
                 key == 'TIER', color)
tf = add_textbox(s, 'TOTAL_COUNT', 2.2, 5.75, 1.8, 0.45)
set_text(tf, '合计', 13, True, DK1)
tf = add_textbox(s, 'TOTAL_SALES', 4.2, 5.75, 2.4, 0.45)
set_text(tf, '0.00', 13, True, DK1)
tf = add_textbox(s, 'CONCLUSION', 0.6, 6.25, 12.1, 0.5)
set_text(tf, '核心结论：____', 15, True, DK1)
cue(s, '讲解提示：头部集中度是分层结论的关键——A 级少数客户贡献大部分销售额。')

# ---- Component 3: List Table (LIST_TABLE, shared by 沉睡/预警/新客/维护) ----
s = new_slide(prs)
header(s, '名单明细', '预警 / 沉睡 / 新客 共用表格')
list_cols = [('COL_NAME', 0.6, 3.4), ('COL_METRIC', 4.2, 2.8), ('COL_STATUS', 7.2, 1.8), ('COL_ACTION', 9.2, 3.5)]
for name, x, w in list_cols:
    tf = add_textbox(s, name, x, 1.7, w, 0.45)
    set_text(tf, {'COL_NAME': '客户名称', 'COL_METRIC': '关键指标', 'COL_STATUS': '状态', 'COL_ACTION': '建议动作'}[name],
             13, True, DK2)
for i in range(1, 6):
    y = 2.25 + (i - 1) * 0.85
    if i % 2 == 1:
        add_rect(s, f'ROW_{i}_BAND', 0.6, y, 12.1, 0.8, fill_color=LT2)
    for name, x, w in list_cols:
        key = name.replace('COL_', '')
        tf = add_textbox(s, f'ROW_{i}_{key}', x, y + 0.18, w, 0.5)
        set_text(tf, '—', 13, False, DK1)
tf = add_textbox(s, 'FOOT_NOTE', 0.6, 6.15, 12.1, 0.5)
set_text(tf, '口径：____（示例：最近 3 个月 vs 前 3 个月）', 11, False, ACC5)
cue(s, '讲解提示：状态列可用主题色区分（预警=警示色、正常=正向色）。')

# ---- Component 4: Customer Card (CUSTOMER_CARD) ----
s = new_slide(prs)
header(s, '客户经营卡片', '单客户画像：层级 / 销售 / 毛利 / 品类 / 趋势 / 建议')
fields = [
    ('FIELD_1', 0.6, 1.75, 5.8, '客户层级', '—'),
    ('FIELD_2', 0.6, 2.55, 5.8, '累计销售额', '—'),
    ('FIELD_3', 0.6, 3.35, 5.8, '累计毛利', '—'),
    ('FIELD_4', 6.9, 1.75, 5.8, '品类结构', '—'),
    ('FIELD_5', 6.9, 2.55, 5.8, '近期趋势', '—'),
    ('FIELD_6', 6.9, 3.35, 5.8, '关注事项', '—'),
]
for prefix, x, y, w, label, value in fields:
    tf = add_textbox(s, prefix + '_LABEL', x, y, 1.9, 0.5)
    set_text(tf, label, 13, False, ACC5)
    tf = add_textbox(s, prefix + '_VALUE', x + 2.0, y, w - 2.0, 0.5)
    set_text(tf, value, 13, True, DK1)
tf = add_textbox(s, 'ADVICE_LABEL', 0.6, 4.5, 2.2, 0.5)
set_text(tf, '下一步建议', 14, True, ACC1)
tf = add_textbox(s, 'ADVICE_TEXT', 2.9, 4.5, 9.8, 1.8)
set_text(tf, '（建议内容：回访节奏 / 促销资源 / 增量品类，标注"供参考"）', 13, False, DK1, line_spacing=1.3)
cue(s, '讲解提示：卡片用于 Top 客户逐一过会；建议必须标注"供参考"。')


def beautify_chart(chart):
    """Remove gridlines and default grey chrome for a clean theme-colored look."""
    try:
        chart.value_axis.has_major_gridlines = False
    except Exception:
        pass
    try:
        chart.category_axis.has_major_gridlines = False
    except Exception:
        pass
    from pptx.oxml.ns import qn
    for area in ('c:chartSpace', 'c:plotArea'):
        el = chart._chartSpace if area == 'c:chartSpace' else chart._chartSpace.find(qn(area))
        if el is None:
            continue
        spPr = el.find(qn('c:spPr'))
        if spPr is None:
            spPr = el.makeelement(qn('c:spPr'), {})
            el.append(spPr)
        else:
            for child in list(spPr):
                spPr.remove(child)
        noFill = spPr.makeelement(qn('a:noFill'), {})
        ln = spPr.makeelement(qn('a:ln'), {})
        lnNoFill = ln.makeelement(qn('a:noFill'), {})
        ln.append(lnNoFill)
        spPr.append(noFill)
        spPr.append(ln)

# ---- Component 5: Pie Chart (CHART_PIE, tier share) ----
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

s = new_slide(prs)
header(s, '客户分层占比', 'A/B/C/D 四级销售额占比（饼图）')
pie_data = CategoryChartData()
pie_data.categories = ['A 级', 'B 级', 'C 级', 'D 级']
pie_data.add_series('销售额占比', (60, 25, 10, 5))   # placeholder data, replaced at use
cf = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(1.6), Inches(1.7), Inches(7.0), Inches(4.4), pie_data)
pie = cf.chart
beautify_chart(pie)
pie.has_legend = True
pie.legend.position = XL_LEGEND_POSITION.RIGHT
pie.legend.include_in_layout = False
pie_colors = [ACC1, ACC6, ACC4, ACC5]
for i, pt in enumerate(pie.plots[0].series[0].points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.theme_color = pie_colors[i]
tf = add_textbox(s, 'CHART_NOTE', 8.9, 2.0, 4.0, 3.8)
set_text(tf, '结论占位：A 级贡献绝大部分销售额，头部集中显著。\n\n数据标签显示各层级占比；替换数据时保持层级顺序 A→B→C→D。', 12, False, ACC5, line_spacing=1.4)
cue(s, '讲解提示：饼图回答"钱从哪里来"；与分层表（CR-02）配合使用。')

# ---- Component 6: Bar Chart (CHART_BAR, monthly trend) ----
s = new_slide(prs)
header(s, '月度销售趋势', '销售额与毛利（柱状图）')
bar_data = CategoryChartData()
bar_data.categories = ['1月', '2月', '3月', '4月', '5月', '6月']
bar_data.add_series('销售额', (100, 90, 60, 50, 75, 60))   # placeholder, replaced at use
bar_data.add_series('毛利', (20, 18, 12, 10, 15, 12))
cf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(1.6), Inches(1.7), Inches(7.6), Inches(4.4), bar_data)
bar = cf.chart
beautify_chart(bar)
bar.has_legend = True
bar.legend.position = XL_LEGEND_POSITION.RIGHT
bar.legend.include_in_layout = False
for series, color in zip(bar.plots[0].series, (ACC1, ACC4)):
    series.format.fill.solid()
    series.format.fill.fore_color.theme_color = color
tf = add_textbox(s, 'CHART_NOTE', 9.5, 2.0, 3.4, 3.8)
set_text(tf, '结论占位：1 月为峰值，Q1 逐月回落，Q2 波动。\n\n柱状图回答"趋势往哪走"。', 12, False, ACC5, line_spacing=1.4)
cue(s, '讲解提示：趋势图回答"走向"，与预警/转化页配合；替换数据时保持 1-6 月顺序。')

prs.save(OUT)
print('components written:', OUT, 'slides:', len(prs.slides))
