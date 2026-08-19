# Page Engine - data-driven slide renderer for the Yuhong component library.
# Replaces one-build-script-per-page: a page is now a JSON spec (blocks), and this
# engine renders any number of specs into one PPTX (based on the Yuhong template master).
#
# Usage:
#   python page-engine.py -o out.pptx spec1.json spec2.json ...
#   python page-engine.py -o out.pptx --template yuhong-template.pptx spec.json
#
# Block types:
#   {"type":"rect",  "name":..., "x":..,"y":..,"w":..,"h":.., "fill":"#505046","line":"#6F6F6A","line_w":0.75}
#   {"type":"text",  "name":..., "x":..,"y":..,"w":..,"h":.., "text":"...","size":14,"bold":false,
#                    "color":"#505046","align":"left|center|right"}
# Coordinates are in inches. Colors are hex (#RRGGBB); the engine converts to RGB.
# Notes (instructor notes) go in "notes" on the spec root.
import argparse
import json

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PALETTE = {
    "ink":   "505046",   # template dk2
    "grey":  "6F6F6A",
    "red":   "B22600",   # deep red (unified; no orange/amber)
    "light": "F7F7F8",
    "white": "FFFFFF",
}


def hex2rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def resolve_color(value):
    return hex2rgb(PALETTE.get(value, value))


def add_block(slide, b):
    x, y, w, h = (Inches(b[k]) for k in ('x', 'y', 'w', 'h'))
    btype = b.get('type')
    if btype == 'rect':
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = resolve_color(b.get('fill', 'white'))
        sh.line.color.rgb = resolve_color(b.get('line', 'grey'))
        sh.line.width = Pt(b.get('line_w', 0.75))
        sh.shadow.inherit = False
        if b.get('name'):
            sh.name = b['name']
    elif btype == 'text':
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = b.get('text', '')
        p.font.size = Pt(b.get('size', 14))
        p.font.bold = b.get('bold', False)
        p.font.name = 'Microsoft YaHei'
        p.font.color.rgb = resolve_color(b.get('color', 'ink'))
        p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}.get(
            b.get('align', 'left'))
        if b.get('name'):
            tb.name = b['name']
    else:
        raise ValueError('unknown block type: %s' % btype)


def render(spec, prs, layout_idx=0):
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    for b in spec.get('blocks', []):
        add_block(slide, b)
    if spec.get('notes'):
        slide.notes_slide.notes_text_frame.text = spec['notes']
    return slide


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('specs', nargs='+')
    ap.add_argument('-o', '--output', required=True)
    ap.add_argument('--template', default=None,
                    help='existing PPTX to use as base (brings its master/layouts)')
    args = ap.parse_args()

    if args.template:
        prs = Presentation(args.template)
    else:
        prs = Presentation()

    for s in args.specs:
        with open(s, encoding='utf-8') as f:
            spec = json.load(f)
        render(spec, prs)

    prs.save(args.output)
    print('rendered %d pages -> %s' % (len(args.specs), args.output))


if __name__ == '__main__':
    main()
